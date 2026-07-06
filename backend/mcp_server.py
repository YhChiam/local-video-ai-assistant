import os
import re
from datetime import datetime
from xml.sax.saxutils import escape
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from ai_engine import LocalAIEngine

ai_engine = LocalAIEngine()


def _is_no_speech_transcript(transcript: str) -> bool:
    text = (transcript or "").strip().lower()
    return text in {"", "(no speech detected)", "no speech detected"}


def _split_sentences(text: str):
    cleaned = re.sub(r"\s+", " ", (text or "").strip())
    if not cleaned:
        return []
    parts = re.split(r"(?<=[.!?])\s+", cleaned)
    return [part.strip(" -") for part in parts if len(part.strip()) >= 20]


def _html_lines(lines: list[str]) -> str:
    return "<br/>".join(escape(line) for line in lines if line)


def _plain_lines(lines: list[str]) -> str:
    return "\n".join(line for line in lines if line)


def _safe_slug(value: str, fallback: str = "video") -> str:
    base_name = os.path.splitext(os.path.basename(value or ""))[0]
    slug = re.sub(r"[^A-Za-z0-9]+", "_", base_name).strip("_").lower()
    return slug or fallback


def _build_output_path(dest_dir: str, prefix: str, source_name: str, extension: str) -> str:
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    slug = _safe_slug(source_name)
    filename = f"{prefix}_{slug}_{timestamp}.{extension}"
    return os.path.join(dest_dir, filename)


def _extract_key_points(transcript: str, vision_text: str, max_points: int = 5):
    points = []
    if not _is_no_speech_transcript(transcript):
        for sentence in _split_sentences(transcript):
            points.append(sentence)
            if len(points) >= max_points:
                break

    if len(points) < max_points and vision_text:
        points.append(f"Visual finding: {vision_text}")

    if not points:
        points = ["No clear spoken key points were detected in this video."]
    return points[:max_points]


def _discussion_summary_lines(transcript: str, vision_text: str, discussion_summary: str = "") -> list[str]:
    lines = []
    if discussion_summary.strip():
        lines.append("Discussion Summary:")
        lines.extend([line for line in re.split(r"\r?\n+", discussion_summary.strip()) if line.strip()])

    lines.append("Transcription Data:")
    lines.append(transcript if not _is_no_speech_transcript(transcript) else "No human speech detected in the provided video audio.")
    lines.append("Vision Analysis Data:")
    lines.append(vision_text)
    return lines

class MCPServerTools:
    """
    Self-developed Model Context Protocol capability mapping layer.
    Keeps application tooling local with no external cloud callbacks.
    """
    @staticmethod
    def call_transcription_tool(video_path: str) -> str:
        return ai_engine.run_whisper_stt(video_path)

    @staticmethod
    def call_vision_tool(video_path: str) -> dict:
        return ai_engine.run_vision_analysis(video_path)

    @staticmethod
    def call_pdf_generation_tool(transcript: str, vision_text: str, dest_dir: str, discussion_summary: str = "", source_name: str = "") -> str:
        os.makedirs(dest_dir, exist_ok=True)
        pdf_path = _build_output_path(dest_dir, "video_analysis_report", source_name, "pdf")
        doc = SimpleDocTemplate(pdf_path, pagesize=letter)
        styles = getSampleStyleSheet()
        transcript_text = transcript if not _is_no_speech_transcript(transcript) else "No human speech detected in the provided video audio."
        story = [
            Paragraph("<b>Video Intelligence Summary Report</b>", styles['Title']),
            Spacer(1, 12),
        ]
        if (discussion_summary or "").strip():
            discussion_lines = [line for line in re.split(r"\r?\n+", discussion_summary.strip()) if line.strip()]
            story.extend([
                Paragraph("<b>Discussion Summary:</b>", styles['Heading2']),
                Paragraph(_html_lines(discussion_lines) or "No prior discussion summary available.", styles['Normal']),
                Spacer(1, 12),
            ])
        story.extend([
            Paragraph("<b>Transcription Data:</b>", styles['Heading2']),
            Paragraph(transcript_text, styles['Normal']),
            Spacer(1, 12),
        ])
        story.extend([
            Paragraph("<b>Vision Analysis Data:</b>", styles['Heading2']),
            Paragraph(vision_text, styles['Normal'])
        ])
        doc.build(story)
        return pdf_path

    @staticmethod
    def call_ppt_generation_tool(transcript: str, vision_text: str, dest_dir: str, discussion_summary: str = "", source_name: str = "") -> str:
        os.makedirs(dest_dir, exist_ok=True)
        ppt_path = _build_output_path(dest_dir, "video_key_points_presentation", source_name, "pptx")
        prs = Presentation()
        transcript_text = transcript if not _is_no_speech_transcript(transcript) else "No speech detected in the source audio."
        discussion_lines = [line for line in re.split(r"\r?\n+", (discussion_summary or "").strip()) if line.strip()]
        if not discussion_lines:
            discussion_lines = ["No prior discussion summary available."]

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Video Key Points Summary"
        title_slide.placeholders[1].text = "Auto-generated from local offline analysis"

        key_points_slide = prs.slides.add_slide(prs.slide_layouts[1])
        key_points_slide.shapes.title.text = "Key Points Discussed"
        key_points = _extract_key_points(transcript_text, vision_text, max_points=5)
        key_points_slide.placeholders[1].text = "\n".join([f"- {point}" for point in key_points])

        details_slide = prs.slides.add_slide(prs.slide_layouts[1])
        details_slide.shapes.title.text = "Analysis Details"
        details_slide.placeholders[1].text = (
            f"Discussion So Far:\n{_plain_lines(discussion_lines)}\n\n"
            f"Transcription Data:\n{transcript_text}\n\n"
            f"Visual Entities:\n{vision_text}"
        )
        
        prs.save(ppt_path)
        return ppt_path